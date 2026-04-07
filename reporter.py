"""
보고서 생성 모듈 — PPT 보고서 전담

CLAUDE.md 규칙:
    - PPT placeholder는 이름으로 접근 (프로그래밍 방식 — shape 직접 추가)
    - 차트 → PNG 먼저 저장 (report/charts/) → PPT 삽입
    - matplotlib: NanumGothic + axes.unicode_minus = False
    - 빈 데이터 → "데이터 없음" 렌더링, 에러 아님
    - logging만 사용, print() 금지
    - 모든 함수에 type hints
"""

import json
import logging
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
import pandas as pd
import seaborn as sns
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────
# 스타일 상수
# ──────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "NanumGothic"
FONT_NAME_BOLD = "NanumGothic Bold"

COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)      # 진한 네이비
COLOR_ACCENT = RGBColor(0xE8, 0x5D, 0x26)        # 주황
COLOR_BG = RGBColor(0xF5, 0xF5, 0xF5)            # 연한 회색 배경
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)

# matplotlib 색상 팔레트
CHART_COLORS = ["#1B3A5C", "#E85D26", "#2ECC71", "#9B59B6",
                "#3498DB", "#E74C3C", "#F39C12", "#1ABC9C",
                "#34495E", "#E67E22"]

NO_DATA_MSG = "데이터 없음"


# ──────────────────────────────────────────────
# 데이터 로딩
# ──────────────────────────────────────────────
@dataclass
class ReportData:
    """보고서 생성에 필요한 모든 데이터를 담는 컨테이너"""
    profile: dict = field(default_factory=dict)
    enriched: pd.DataFrame = field(default_factory=pd.DataFrame)
    format_stats: pd.DataFrame = field(default_factory=pd.DataFrame)
    categories: dict | None = None
    caption_style: dict | None = None
    sentiment: dict | None = None
    audience_age: dict | None = None
    visual: dict | None = None
    top_posts: dict | None = None
    insights: dict | None = None
    images_dir: Path = Path()
    charts_dir: Path = Path()


def _load_json(path: Path) -> dict | None:
    """JSON 파일 안전 로딩 — 없으면 None 반환"""
    if not path.exists():
        logger.debug("JSON 파일 없음 (건너뜀): %s", path)
        return None
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError) as e:
        logger.warning("JSON 로딩 실패: %s — %s", path, e)
        return None


def _load_csv(path: Path) -> pd.DataFrame | None:
    """CSV 파일 안전 로딩 — 없으면 None 반환"""
    if not path.exists():
        logger.debug("CSV 파일 없음 (건너뜀): %s", path)
        return None
    try:
        return pd.read_csv(path)
    except (pd.errors.EmptyDataError, OSError) as e:
        logger.warning("CSV 로딩 실패: %s — %s", path, e)
        return None


def _load_report_data(data_dir: Path) -> ReportData:
    """보고서에 필요한 모든 데이터를 로딩"""
    raw_dir = data_dir / "raw"
    analysis_dir = data_dir / "analysis"
    charts_dir = data_dir / "report" / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)

    profile = _load_json(raw_dir / "profile.json") or {}
    enriched = _load_csv(analysis_dir / "posts_enriched.csv")
    if enriched is None:
        enriched = _load_csv(raw_dir / "posts.csv")
    if enriched is None:
        enriched = pd.DataFrame()

    format_stats = _load_csv(analysis_dir / "format_stats.csv")
    if format_stats is None:
        format_stats = pd.DataFrame()

    return ReportData(
        profile=profile,
        enriched=enriched,
        format_stats=format_stats,
        categories=_load_json(analysis_dir / "categories.json"),
        caption_style=_load_json(analysis_dir / "caption_style.json"),
        sentiment=_load_json(analysis_dir / "sentiment.json"),
        audience_age=_load_json(analysis_dir / "audience_age.json"),
        visual=_load_json(analysis_dir / "visual.json"),
        top_posts=_load_json(analysis_dir / "top_posts_analysis.json"),
        insights=_load_json(analysis_dir / "insights.json"),
        images_dir=raw_dir / "images",
        charts_dir=charts_dir,
    )


def _get_insight_section(insights: dict | None, section: str) -> str:
    """insights에서 특정 섹션 내러티브 추출 — 없으면 NO_DATA_MSG"""
    if not insights:
        return NO_DATA_MSG
    sections = insights.get("sections", {})
    return sections.get(section, NO_DATA_MSG) or NO_DATA_MSG


# ──────────────────────────────────────────────
# matplotlib 설정
# ──────────────────────────────────────────────
def _setup_matplotlib() -> None:
    """matplotlib 한국어 폰트 및 스타일 설정"""
    # matplotlib 폰트 캐시 갱신
    fm._load_fontmanager(try_read_cache=False)

    plt.rcParams["font.family"] = FONT_NAME
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.dpi"] = 150
    plt.rcParams["savefig.dpi"] = 150
    plt.rcParams["savefig.bbox"] = "tight"
    plt.rcParams["figure.facecolor"] = "white"

    # 폰트 존재 확인
    font_paths = fm.findSystemFonts()
    nanum_found = any("NanumGothic" in p for p in font_paths)
    if not nanum_found:
        logger.warning("NanumGothic 폰트를 찾을 수 없습니다. 기본 폰트를 사용합니다.")
    else:
        logger.debug("NanumGothic 폰트 확인됨")


# ──────────────────────────────────────────────
# 차트 생성 함수
# ──────────────────────────────────────────────
def _chart_format_comparison(format_stats: pd.DataFrame, charts_dir: Path) -> Path | None:
    """포맷별 성과 비교 바 차트"""
    if format_stats.empty:
        return None

    output = charts_dir / "format_comparison.png"
    fig, ax = plt.subplots(figsize=(10, 5))

    formats = format_stats["typename"].tolist()
    # 한국어 포맷명 매핑
    format_labels = {
        "GraphImage": "이미지", "GraphVideo": "릴스",
        "GraphSidecar": "캐러셀",
    }
    labels = [format_labels.get(f, f) for f in formats]

    x = np.arange(len(labels))
    width = 0.25

    likes = format_stats.get("avg_likes", pd.Series([0] * len(formats))).fillna(0)
    comments = format_stats.get("avg_comments", pd.Series([0] * len(formats))).fillna(0)
    saves = format_stats.get("avg_est_saves", pd.Series([0] * len(formats))).fillna(0)

    ax.bar(x - width, likes, width, label="평균 좋아요", color=CHART_COLORS[0])
    ax.bar(x, comments, width, label="평균 댓글", color=CHART_COLORS[1])
    ax.bar(x + width, saves, width, label="추정 저장", color=CHART_COLORS[2])

    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylabel("수")
    ax.set_title("포맷별 평균 성과 비교")
    ax.legend()

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _chart_category_pie(
    enriched: pd.DataFrame, categories: dict | None, charts_dir: Path
) -> Path | None:
    """콘텐츠 카테고리 분포 파이 차트"""
    if categories is None:
        return None

    classifications = categories.get("classifications", [])
    if not classifications:
        return None

    output = charts_dir / "category_pie.png"

    cat_counts = Counter(c.get("category_code", "기타") for c in classifications)
    # category_name 매핑
    code_to_name = {}
    for c in classifications:
        code = c.get("category_code", "")
        name = c.get("category_name", code)
        if code not in code_to_name:
            code_to_name[code] = name

    labels = [code_to_name.get(k, k) for k in cat_counts.keys()]
    sizes = list(cat_counts.values())
    colors = CHART_COLORS[:len(labels)]

    fig, ax = plt.subplots(figsize=(8, 8))
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, autopct="%1.1f%%",
        colors=colors, startangle=90, pctdistance=0.85,
    )
    for text in autotexts:
        text.set_fontsize(9)
    ax.set_title("콘텐츠 카테고리 분포")

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _chart_posting_frequency(enriched: pd.DataFrame, charts_dir: Path) -> Path | None:
    """게시 빈도 추이 라인 차트"""
    if enriched.empty or "date_utc" not in enriched.columns:
        return None

    output = charts_dir / "posting_frequency.png"

    df = enriched.copy()
    df["date_utc"] = pd.to_datetime(df["date_utc"], errors="coerce")
    df = df.dropna(subset=["date_utc"])
    if df.empty:
        return None

    # 주간 빈도
    df["week"] = df["date_utc"].dt.to_period("W").dt.start_time
    weekly = df.groupby("week").size().reset_index(name="count")

    fig, ax = plt.subplots(figsize=(10, 4))
    ax.plot(weekly["week"], weekly["count"], marker="o",
            color=CHART_COLORS[0], linewidth=2, markersize=4)
    ax.fill_between(weekly["week"], weekly["count"], alpha=0.1, color=CHART_COLORS[0])
    ax.set_xlabel("기간")
    ax.set_ylabel("게시물 수")
    ax.set_title("주간 게시 빈도 추이")
    fig.autofmt_xdate()

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _chart_day_of_week(enriched: pd.DataFrame, charts_dir: Path) -> Path | None:
    """요일별 평균 성과 바 차트"""
    if enriched.empty or "date_utc" not in enriched.columns:
        return None

    output = charts_dir / "day_of_week.png"

    df = enriched.copy()
    df["date_utc"] = pd.to_datetime(df["date_utc"], errors="coerce")
    df = df.dropna(subset=["date_utc"])
    if df.empty:
        return None

    df["dow"] = df["date_utc"].dt.dayofweek  # 0=Mon
    day_labels = ["월", "화", "수", "목", "금", "토", "일"]

    dow_stats = df.groupby("dow").agg(
        count=("likes", "size"),
        avg_likes=("likes", "mean"),
    ).reindex(range(7), fill_value=0)

    fig, ax1 = plt.subplots(figsize=(8, 5))
    ax1.bar(range(7), dow_stats["count"], color=CHART_COLORS[0], alpha=0.7, label="게시물 수")
    ax1.set_ylabel("게시물 수")
    ax1.set_xticks(range(7))
    ax1.set_xticklabels(day_labels)

    ax2 = ax1.twinx()
    ax2.plot(range(7), dow_stats["avg_likes"], color=CHART_COLORS[1],
             marker="o", linewidth=2, label="평균 좋아요")
    ax2.set_ylabel("평균 좋아요")

    ax1.set_title("요일별 게시 빈도 및 성과")
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper right")

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _chart_time_heatmap(enriched: pd.DataFrame, charts_dir: Path) -> Path | None:
    """시간대별 게시물 히트맵 (요일 x 시간)"""
    if enriched.empty or "date_utc" not in enriched.columns:
        return None

    output = charts_dir / "time_heatmap.png"

    df = enriched.copy()
    df["date_utc"] = pd.to_datetime(df["date_utc"], errors="coerce")
    df = df.dropna(subset=["date_utc"])
    if df.empty:
        return None

    # KST 변환 (+9시간)
    df["kst"] = df["date_utc"] + pd.Timedelta(hours=9)
    df["dow"] = df["kst"].dt.dayofweek
    df["hour"] = df["kst"].dt.hour

    pivot = df.pivot_table(
        index="dow", columns="hour", values="likes",
        aggfunc="count", fill_value=0,
    )
    # 전체 시간대 보장
    pivot = pivot.reindex(index=range(7), columns=range(24), fill_value=0)

    day_labels = ["월", "화", "수", "목", "금", "토", "일"]

    fig, ax = plt.subplots(figsize=(12, 4))
    sns.heatmap(
        pivot, ax=ax, cmap="YlOrRd", annot=True, fmt="g",
        xticklabels=[f"{h}시" for h in range(24)],
        yticklabels=day_labels,
        cbar_kws={"label": "게시물 수"},
    )
    ax.set_title("요일 × 시간대 게시 분포 (KST)")
    ax.set_xlabel("시간")
    ax.set_ylabel("요일")

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _chart_sentiment_pie(sentiment: dict | None, charts_dir: Path) -> Path | None:
    """댓글 감성 분포 파이 차트"""
    if sentiment is None:
        return None

    dist = sentiment.get("overall_sentiment_distribution", {})
    if not dist:
        return None

    output = charts_dir / "sentiment_pie.png"

    labels_map = {"positive": "긍정", "neutral": "중립", "negative": "부정"}
    colors_map = {"positive": "#2ECC71", "neutral": "#95A5A6", "negative": "#E74C3C"}

    labels = []
    sizes = []
    colors = []
    for key in ["positive", "neutral", "negative"]:
        val = dist.get(key, 0)
        if val > 0:
            labels.append(labels_map.get(key, key))
            sizes.append(val)
            colors.append(colors_map.get(key, "#999999"))

    if not sizes:
        return None

    fig, ax = plt.subplots(figsize=(6, 6))
    ax.pie(sizes, labels=labels, autopct="%1.1f%%", colors=colors, startangle=90)
    ax.set_title("댓글 감성 분포")

    fig.savefig(output)
    plt.close(fig)
    logger.info("차트 저장: %s", output)
    return output


def _create_image_grid(
    images_dir: Path, shortcodes: list[str], charts_dir: Path
) -> Path | None:
    """게시물 이미지 그리드 (3x3) 생성"""
    if not images_dir.exists():
        return None

    # 이미지 찾기
    available = []
    for sc in shortcodes[:9]:
        img_path = images_dir / f"{sc}.jpg"
        if img_path.exists():
            available.append(img_path)
    if not available:
        return None

    output = charts_dir / "image_grid.png"
    cell_size = 300
    cols = min(3, len(available))
    rows = (len(available) + cols - 1) // cols

    grid = Image.new("RGB", (cols * cell_size, rows * cell_size), (245, 245, 245))
    for i, img_path in enumerate(available):
        try:
            img = Image.open(img_path)
            img.thumbnail((cell_size, cell_size))
            r, c = divmod(i, cols)
            # 센터 정렬
            x_offset = c * cell_size + (cell_size - img.width) // 2
            y_offset = r * cell_size + (cell_size - img.height) // 2
            grid.paste(img, (x_offset, y_offset))
        except OSError:
            continue

    grid.save(output)
    logger.info("이미지 그리드 저장: %s", output)
    return output


# ──────────────────────────────────────────────
# PPT 헬퍼 함수
# ──────────────────────────────────────────────
def _set_text_style(
    tf, text: str, font_size: int = 14,
    bold: bool = False, color: RGBColor = COLOR_TEXT,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """텍스트 프레임에 스타일 적용된 텍스트 설정"""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment


def _add_textbox(
    slide, left: float, top: float, width: float, height: float,
    text: str, font_size: int = 14, bold: bool = False,
    color: RGBColor = COLOR_TEXT, alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """슬라이드에 텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text_style(tf, text, font_size, bold, color, alignment)


def _add_multiline_textbox(
    slide, left: float, top: float, width: float, height: float,
    lines: list[tuple[str, int, bool, RGBColor]],
) -> None:
    """여러 줄 텍스트 박스 추가

    Args:
        lines: [(text, font_size, bold, color), ...]
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)


def _add_section_header(slide, title: str) -> None:
    """섹션 제목 바 추가 (상단 배너)"""
    # 배경 바
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.0)  # 1 = Rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    _add_textbox(slide, 0.5, 0.15, 12, 0.7, title,
                 font_size=24, bold=True, color=COLOR_WHITE)


def _add_image_safe(slide, img_path: Path, left: float, top: float,
                    width: float | None = None, height: float | None = None) -> bool:
    """이미지를 안전하게 삽입 — 없으면 False 반환"""
    if not img_path or not img_path.exists():
        return False
    try:
        kwargs = {"image_file": str(img_path),
                  "left": Inches(left), "top": Inches(top)}
        if width:
            kwargs["width"] = Inches(width)
        if height:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(**kwargs)
        return True
    except Exception as e:
        logger.warning("이미지 삽입 실패: %s — %s", img_path, e)
        return False


def _add_table(
    slide, left: float, top: float, width: float,
    headers: list[str], rows: list[list[str]],
    col_widths: list[float] | None = None,
) -> None:
    """슬라이드에 테이블 추가"""
    n_rows = len(rows) + 1  # 헤더 포함
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(0.4 * n_rows),
    )
    table = table_shape.table

    # 열 너비 설정
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # 헤더 행
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = FONT_NAME
                paragraph.font.color.rgb = COLOR_TEXT


# ──────────────────────────────────────────────
# 슬라이드 빌더 (섹션별)
# ──────────────────────────────────────────────
def _build_cover(prs: Presentation, data: ReportData) -> None:
    """Slide 1: 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 배경색
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    username = data.profile.get("username", "채널명")
    _add_textbox(slide, 1, 1.5, 11, 1.5,
                 "인스타그램 채널 분석 보고서",
                 font_size=36, bold=True, color=COLOR_WHITE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, 1, 3.2, 11, 1,
                 f"@{username}",
                 font_size=28, color=COLOR_ACCENT,
                 alignment=PP_ALIGN.CENTER)

    # 분석 기간
    if not data.enriched.empty and "date_utc" in data.enriched.columns:
        dates = pd.to_datetime(data.enriched["date_utc"], errors="coerce").dropna()
        if not dates.empty:
            period = f"{dates.min().strftime('%Y.%m.%d')} ~ {dates.max().strftime('%Y.%m.%d')}"
        else:
            period = ""
    else:
        period = ""

    subtitle = f"분석 기간: {period}" if period else ""
    subtitle += f"\n생성일: {datetime.now().strftime('%Y.%m.%d')}"
    _add_textbox(slide, 1, 4.5, 11, 1, subtitle,
                 font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB),
                 alignment=PP_ALIGN.CENTER)


def _build_executive_summary(prs: Presentation, data: ReportData) -> None:
    """Slide 2: §1 분석 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§1 분석 개요 — Executive Summary")

    summary = NO_DATA_MSG
    if data.insights:
        summary = data.insights.get("executive_summary", NO_DATA_MSG) or NO_DATA_MSG

    _add_textbox(slide, 0.8, 1.3, 11.5, 5.5, summary, font_size=16)

    # 핵심 지표 콜아웃
    if data.insights:
        callouts = data.insights.get("key_metrics_callout", [])
        if callouts:
            lines = [(f"  {c}", 13, False, COLOR_ACCENT) for c in callouts[:3]]
            _add_multiline_textbox(slide, 0.8, 5.0, 11.5, 2.0, lines)


def _build_channel_profile(prs: Presentation, data: ReportData) -> None:
    """Slide 3: §2 채널 프로필"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§2 채널 프로필 스냅샷")

    p = data.profile
    followers = f"{p.get('followers', 0):,}"
    following = f"{p.get('following', 0):,}"
    media_count = f"{p.get('mediacount', 0):,}"
    bio = p.get("biography", "")
    verified = "있음" if p.get("is_verified", False) else "없음"
    category = p.get("category", "") or p.get("category_name", "") or ""
    link = p.get("external_url", "") or ""

    info_lines = [
        (f"계정명: @{p.get('username', '')}", 14, True, COLOR_TEXT),
        (f"팔로워: {followers}  |  팔로잉: {following}  |  게시물: {media_count}", 13, False, COLOR_TEXT),
        (f"인증 배지: {verified}  |  카테고리: {category or '미분류'}", 13, False, COLOR_TEXT),
        (f"외부 링크: {link or '없음'}", 12, False, COLOR_LIGHT_TEXT),
        ("", 10, False, COLOR_TEXT),
        (f"Bio: {bio}", 12, False, COLOR_TEXT),
    ]
    _add_multiline_textbox(slide, 0.8, 1.3, 7, 3.5, info_lines)

    # 내러티브
    narrative = _get_insight_section(data.insights, "channel_profile")
    _add_textbox(slide, 0.8, 5.0, 11.5, 2.0, narrative, font_size=12, color=COLOR_LIGHT_TEXT)


def _build_audience(prs: Presentation, data: ReportData) -> None:
    """Slides 4-5: §3 오디언스 추정"""
    # Slide 4: 데모그래픽
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§3 오디언스 추정 분석")

    if data.audience_age:
        age_dist = data.audience_age.get("age_distribution", {})
        if age_dist:
            headers = ["연령대", "추정 비율"]
            rows = [[k, f"{v}%"] for k, v in age_dist.items()]
            _add_table(slide, 0.8, 1.5, 5, headers, rows)
        else:
            _add_textbox(slide, 0.8, 1.5, 11.5, 1, NO_DATA_MSG, font_size=14)
    else:
        _add_textbox(slide, 0.8, 1.5, 11.5, 1,
                     "오디언스 연령대 추정 데이터 없음 (댓글 수집 또는 AI 분석 필요)",
                     font_size=14, color=COLOR_LIGHT_TEXT)

    narrative = _get_insight_section(data.insights, "audience")
    _add_textbox(slide, 0.8, 5.0, 11.5, 2.0, narrative, font_size=12, color=COLOR_LIGHT_TEXT)

    # Slide 5: 활동 시간대 히트맵
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§3 추정 활동 시간대")

    heatmap_path = _chart_time_heatmap(data.enriched, data.charts_dir)
    if heatmap_path:
        _add_image_safe(slide, heatmap_path, 0.5, 1.2, width=12)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)


def _build_content_strategy(prs: Presentation, data: ReportData) -> None:
    """Slides 6-8: §4 콘텐츠 전략"""
    format_labels = {"GraphImage": "이미지", "GraphVideo": "릴스", "GraphSidecar": "캐러셀"}

    # Slide 6: 포맷별 성과
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§4 포맷별 성과 비교")

    fmt_chart = _chart_format_comparison(data.format_stats, data.charts_dir)
    if not data.format_stats.empty:
        headers = ["포맷", "게시물 수", "비율(%)", "평균 좋아요", "평균 댓글",
                   "추정 저장", "추정 공유", "추정 인게이지먼트율"]
        rows = []
        for _, r in data.format_stats.iterrows():
            rows.append([
                format_labels.get(r.get("typename", ""), r.get("typename", "")),
                f"{int(r.get('count', 0)):,}",
                f"{r.get('ratio_pct', 0):.1f}",
                f"{r.get('avg_likes', 0):,.0f}",
                f"{r.get('avg_comments', 0):,.0f}",
                f"{r.get('avg_est_saves', 0):,.0f}",
                f"{r.get('avg_est_shares', 0):,.0f}",
                f"{r.get('avg_est_eng_rate', 0):.2f}%",
            ])
        _add_table(slide, 0.3, 1.3, 12.5, headers, rows)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)

    if fmt_chart:
        _add_image_safe(slide, fmt_chart, 3, 4.2, width=7)

    # Slide 7: 카테고리 분석
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§4 콘텐츠 카테고리 분류 (AI)")

    cat_chart = _chart_category_pie(data.enriched, data.categories, data.charts_dir)
    if data.categories and not data.enriched.empty:
        # 카테고리별 성과 집계
        classifications = data.categories.get("classifications", [])
        sc_to_cat = {c["shortcode"]: (c.get("category_code", ""), c.get("category_name", ""))
                     for c in classifications if "shortcode" in c}

        df = data.enriched.copy()
        df["cat_code"] = df["shortcode"].map(lambda s: sc_to_cat.get(s, ("", ""))[0])
        df["cat_name"] = df["shortcode"].map(lambda s: sc_to_cat.get(s, ("", ""))[1])

        cat_stats = df.groupby(["cat_code", "cat_name"]).agg(
            count=("likes", "size"),
            avg_likes=("likes", "mean"),
            avg_comments=("comments", "mean"),
        ).reset_index()
        cat_stats = cat_stats.sort_values("count", ascending=False)
        total = cat_stats["count"].sum()

        headers = ["카테고리", "게시물 수", "비율(%)", "평균 좋아요", "평균 댓글"]
        rows = []
        for _, r in cat_stats.iterrows():
            name = r["cat_name"] if r["cat_name"] else r["cat_code"]
            pct = (r["count"] / total * 100) if total > 0 else 0
            rows.append([
                name or "기타",
                f"{int(r['count'])}",
                f"{pct:.1f}",
                f"{r['avg_likes']:,.0f}",
                f"{r['avg_comments']:,.0f}",
            ])
        _add_table(slide, 0.3, 1.3, 6.5, headers, rows[:10])

        if cat_chart:
            _add_image_safe(slide, cat_chart, 7.5, 1.3, width=5)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1,
                     "카테고리 분석 데이터 없음 (AI 분석 필요)", font_size=14, color=COLOR_LIGHT_TEXT)

    narrative = _get_insight_section(data.insights, "content_strategy")
    _add_textbox(slide, 0.3, 5.5, 12.5, 1.5, narrative, font_size=11, color=COLOR_LIGHT_TEXT)

    # Slide 8: 캡션 스타일
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§4 캡션 스타일 분석")

    if data.caption_style:
        cs = data.caption_style
        info_lines = [
            (f"톤앤매너: {cs.get('tone', NO_DATA_MSG)}", 14, True, COLOR_TEXT),
            (f"평균 캡션 길이: {cs.get('avg_length', '?')}자", 13, False, COLOR_TEXT),
        ]
        # CTA 유형
        cta_types = cs.get("cta_types", [])
        if cta_types:
            if isinstance(cta_types, list):
                cta_str = ", ".join(str(c) for c in cta_types[:5])
            else:
                cta_str = str(cta_types)
            info_lines.append((f"CTA 유형: {cta_str}", 13, False, COLOR_TEXT))

        # 이모지 패턴
        emoji_patterns = cs.get("emoji_patterns", "")
        if emoji_patterns:
            info_lines.append((f"이모지 패턴: {emoji_patterns}", 13, False, COLOR_TEXT))

        _add_multiline_textbox(slide, 0.8, 1.3, 11.5, 3.5, info_lines)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)

    narrative = _get_insight_section(data.insights, "caption_style")
    _add_textbox(slide, 0.8, 5.0, 11.5, 2.0, narrative, font_size=12, color=COLOR_LIGHT_TEXT)


def _build_visual_tone(prs: Presentation, data: ReportData) -> None:
    """Slides 9-10: §5 비주얼 톤"""
    # Slide 9: 피드 심미성
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§5 비주얼 톤 & 크리에이티브 분석")

    if data.visual:
        v = data.visual
        info_lines = [
            (f"피드 통일감 점수: {v.get('feed_coherence_score', '?')}/10", 16, True, COLOR_TEXT),
        ]
        palette = v.get("color_palette", [])
        if palette:
            info_lines.append((f"주조 색감: {', '.join(str(c) for c in palette[:5])}", 13, False, COLOR_TEXT))

        styles = v.get("design_styles", [])
        if styles:
            info_lines.append((f"이미지 스타일: {', '.join(str(s) for s in styles[:5])}", 13, False, COLOR_TEXT))

        style_dist = v.get("design_style_distribution", {})
        if style_dist:
            dist_str = " | ".join(f"{k}: {v2}%" for k, v2 in style_dist.items())
            info_lines.append((f"스타일 분포: {dist_str}", 12, False, COLOR_LIGHT_TEXT))

        _add_multiline_textbox(slide, 0.8, 1.3, 11.5, 3.5, info_lines)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1,
                     "비주얼 분석 데이터 없음 (Vision 분석 필요)",
                     font_size=14, color=COLOR_LIGHT_TEXT)

    narrative = _get_insight_section(data.insights, "visual_tone")
    _add_textbox(slide, 0.8, 5.0, 11.5, 2.0, narrative, font_size=12, color=COLOR_LIGHT_TEXT)

    # Slide 10: 이미지 그리드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§5 최근 게시물 피드 미리보기")

    if not data.enriched.empty and "shortcode" in data.enriched.columns:
        shortcodes = data.enriched["shortcode"].tolist()[:9]
        grid_path = _create_image_grid(data.images_dir, shortcodes, data.charts_dir)
        if grid_path:
            _add_image_safe(slide, grid_path, 2.5, 1.3, width=8)
        else:
            _add_textbox(slide, 0.8, 3, 11.5, 1,
                         "게시물 이미지 파일 없음", font_size=14, color=COLOR_LIGHT_TEXT)
    else:
        _add_textbox(slide, 0.8, 3, 11.5, 1, NO_DATA_MSG, font_size=14)


def _build_top_posts(prs: Presentation, data: ReportData) -> None:
    """Slides 11-13: §6 TOP 인기 게시물"""
    format_labels = {"GraphImage": "이미지", "GraphVideo": "릴스", "GraphSidecar": "캐러셀"}

    # Slide 11: TOP 10 리스트
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§6 좋아요 기준 TOP 10 게시물")

    if not data.enriched.empty and "likes" in data.enriched.columns:
        top10 = data.enriched.nlargest(10, "likes")

        # 카테고리 매핑
        sc_to_cat = {}
        if data.categories:
            for c in data.categories.get("classifications", []):
                sc_to_cat[c.get("shortcode", "")] = c.get("category_name", "")

        headers = ["순위", "게시일", "포맷", "카테고리", "좋아요", "댓글", "추정 인게이지먼트율"]
        rows = []
        for rank, (_, r) in enumerate(top10.iterrows(), 1):
            date_str = ""
            if "date_utc" in r and pd.notna(r["date_utc"]):
                try:
                    date_str = pd.to_datetime(r["date_utc"]).strftime("%Y.%m.%d")
                except (ValueError, TypeError):
                    date_str = str(r["date_utc"])[:10]
            fmt = format_labels.get(r.get("typename", ""), r.get("typename", ""))
            cat = sc_to_cat.get(r.get("shortcode", ""), "")
            eng_rate = r.get("est_full_eng_rate", r.get("est_eng_rate", 0))
            rows.append([
                str(rank), date_str, fmt, cat or "-",
                f"{int(r.get('likes', 0)):,}",
                f"{int(r.get('comments', 0)):,}",
                f"{eng_rate:.2f}%" if eng_rate else "-",
            ])
        _add_table(slide, 0.3, 1.3, 12.5, headers, rows)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)

    # Slide 12: 인기 게시물 공통점
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§6 인기 게시물 공통점 분석")

    if data.top_posts:
        tp = data.top_posts
        info_lines = []
        for key, label in [
            ("success_factors", "성공 요인"),
            ("format_distribution", "포맷 분포"),
            ("category_patterns", "카테고리 패턴"),
            ("timing_patterns", "타이밍 패턴"),
        ]:
            val = tp.get(key)
            if val:
                if isinstance(val, list):
                    val_str = ", ".join(str(v) for v in val[:5])
                elif isinstance(val, dict):
                    val_str = " | ".join(f"{k}: {v}" for k, v in val.items())
                else:
                    val_str = str(val)
                info_lines.append((f"{label}: {val_str}", 12, False, COLOR_TEXT))
        if info_lines:
            _add_multiline_textbox(slide, 0.8, 1.3, 11.5, 4, info_lines)
        else:
            _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1,
                     "인기 게시물 분석 데이터 없음 (AI 분석 필요)",
                     font_size=14, color=COLOR_LIGHT_TEXT)

    narrative = _get_insight_section(data.insights, "top_posts")
    _add_textbox(slide, 0.8, 5.0, 11.5, 2.0, narrative, font_size=12, color=COLOR_LIGHT_TEXT)

    # Slide 13: 감성 분석
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§6 댓글 감성 분석")

    sentiment_chart = _chart_sentiment_pie(data.sentiment, data.charts_dir)
    if sentiment_chart:
        _add_image_safe(slide, sentiment_chart, 3.5, 1.2, width=6)
    elif data.sentiment:
        dist = data.sentiment.get("overall_sentiment_distribution", {})
        if dist:
            _add_textbox(slide, 0.8, 2, 11.5, 1,
                         f"긍정: {dist.get('positive', 0)}% | 중립: {dist.get('neutral', 0)}% | 부정: {dist.get('negative', 0)}%",
                         font_size=16)
        else:
            _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1,
                     "감성 분석 데이터 없음 (댓글 수집 및 AI 분석 필요)",
                     font_size=14, color=COLOR_LIGHT_TEXT)


def _build_posting_timing(prs: Presentation, data: ReportData) -> None:
    """Slides 14-15: §7 게시 빈도 & 타이밍"""
    # Slide 14: 빈도 추이
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§7 게시 빈도 추이")

    freq_chart = _chart_posting_frequency(data.enriched, data.charts_dir)
    if freq_chart:
        _add_image_safe(slide, freq_chart, 0.5, 1.2, width=12)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)

    # 주간 평균 게시 횟수
    if not data.enriched.empty and "date_utc" in data.enriched.columns:
        dates = pd.to_datetime(data.enriched["date_utc"], errors="coerce").dropna()
        if not dates.empty:
            weeks = (dates.max() - dates.min()).days / 7
            avg_per_week = len(dates) / max(weeks, 1)
            _add_textbox(slide, 0.8, 5.5, 5, 0.5,
                         f"주간 평균 게시 횟수: {avg_per_week:.1f}회",
                         font_size=13, bold=True)

    narrative = _get_insight_section(data.insights, "posting_timing")
    _add_textbox(slide, 0.8, 6.2, 11.5, 1, narrative, font_size=11, color=COLOR_LIGHT_TEXT)

    # Slide 15: 요일별 분석
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§7 요일별 / 시간대별 성과")

    dow_chart = _chart_day_of_week(data.enriched, data.charts_dir)
    if dow_chart:
        _add_image_safe(slide, dow_chart, 0.5, 1.2, width=12)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)


def _build_hashtags(prs: Presentation, data: ReportData) -> None:
    """Slides 16-17: §8 해시태그 전략"""
    # Slide 16: TOP 20 해시태그
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§8 해시태그 TOP 20")

    if not data.enriched.empty and "hashtags" in data.enriched.columns:
        # 해시태그 빈도 집계
        all_tags: Counter = Counter()
        tag_likes: dict[str, list[float]] = {}

        for _, row in data.enriched.iterrows():
            tags_str = str(row.get("hashtags", ""))
            if not tags_str or tags_str == "nan":
                continue
            tags = [t.strip().strip("#") for t in tags_str.split(",") if t.strip()]
            likes_val = float(row.get("likes", 0))
            for tag in tags:
                if tag:
                    all_tags[tag] += 1
                    tag_likes.setdefault(tag, []).append(likes_val)

        top20 = all_tags.most_common(20)
        if top20:
            avg_per_post = sum(all_tags.values()) / len(data.enriched) if len(data.enriched) > 0 else 0
            _add_textbox(slide, 0.8, 1.2, 5, 0.4,
                         f"게시물당 평균 해시태그 수: {avg_per_post:.1f}개",
                         font_size=12, bold=True)

            headers = ["순위", "해시태그", "사용 횟수", "해당 게시물 평균 좋아요"]
            rows = []
            for i, (tag, cnt) in enumerate(top20, 1):
                avg_like = np.mean(tag_likes.get(tag, [0]))
                rows.append([str(i), f"#{tag}", str(cnt), f"{avg_like:,.0f}"])
            _add_table(slide, 0.3, 1.8, 12.5, headers, rows)
        else:
            _add_textbox(slide, 0.8, 2.5, 11.5, 1, "해시태그 데이터 없음", font_size=14)
    else:
        _add_textbox(slide, 0.8, 2.5, 11.5, 1, NO_DATA_MSG, font_size=14)

    # Slide 17: 해시태그 전략 내러티브
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "§8 해시태그 전략 인사이트")

    narrative = _get_insight_section(data.insights, "hashtag_strategy")
    _add_textbox(slide, 0.8, 1.5, 11.5, 5, narrative, font_size=14)


def _build_closing(prs: Presentation, data: ReportData) -> None:
    """Slide 18: 마무리 — 핵심 인사이트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    _add_textbox(slide, 1, 0.8, 11, 1,
                 "핵심 인사이트 & 추천 사항",
                 font_size=28, bold=True, color=COLOR_WHITE,
                 alignment=PP_ALIGN.CENTER)

    recommendations = _get_insight_section(data.insights, "recommendations")
    _add_textbox(slide, 1.5, 2.2, 10, 4.5, recommendations,
                 font_size=15, color=COLOR_WHITE)


def _build_appendix(prs: Presentation) -> None:
    """Slide 19: 부록 — 면책 조항"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, "부록 — 데이터 한계 및 면책 사항")

    disclaimer = (
        "1. 본 보고서는 공개 데이터 + 업계 벤치마크 기반 추정치로 작성되었습니다.\n\n"
        "2. 🔶 추정 표기된 지표의 계수는 글로벌 평균 데이터 기반이며, "
        "한국 시장/특정 업종 특화 계수는 아닙니다.\n\n"
        "3. 저장/공유 계수 출처: Socialinsider 2026, Sprout Social 2025, "
        "Statista/Metricool 2024, Dash Social 2025\n\n"
        "4. 도달(Views)과 순 도달(Reach)은 다른 지표입니다. "
        "순 도달은 조회 수보다 낮습니다.\n\n"
        "5. Instagram 알고리즘은 지속적으로 변하며, "
        "분석 시점에 따라 데이터의 유효성이 달라질 수 있습니다.\n\n"
        "6. 프로필 방문 / 웹사이트 클릭은 신뢰할 수 있는 추정 방법이 없어 제외하였습니다."
    )
    _add_textbox(slide, 0.8, 1.5, 11.5, 5.5, disclaimer, font_size=12, color=COLOR_LIGHT_TEXT)


# ──────────────────────────────────────────────
# 메인 진입점
# ──────────────────────────────────────────────
def generate_report(channel: str, data_dir: Path) -> Path:
    """보고서 생성 메인 함수

    Args:
        channel: 채널명 (@없이)
        data_dir: data/{channel}/ 경로

    Returns:
        생성된 report.pptx 경로
    """
    logger.info("보고서 생성 시작 — @%s", channel)

    # matplotlib 한국어 폰트 설정
    _setup_matplotlib()

    # 데이터 로딩
    data = _load_report_data(data_dir)
    logger.info(
        "데이터 로딩 완료 — 게시물 %d개, AI분석: %s",
        len(data.enriched),
        "있음" if data.insights else "없음",
    )

    # PPT 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 각 섹션 빌드
    _build_cover(prs, data)
    logger.debug("슬라이드 1: 표지 생성 완료")

    _build_executive_summary(prs, data)
    logger.debug("슬라이드 2: 분석 개요 생성 완료")

    _build_channel_profile(prs, data)
    logger.debug("슬라이드 3: 채널 프로필 생성 완료")

    _build_audience(prs, data)
    logger.debug("슬라이드 4-5: 오디언스 추정 생성 완료")

    _build_content_strategy(prs, data)
    logger.debug("슬라이드 6-8: 콘텐츠 전략 생성 완료")

    _build_visual_tone(prs, data)
    logger.debug("슬라이드 9-10: 비주얼 톤 생성 완료")

    _build_top_posts(prs, data)
    logger.debug("슬라이드 11-13: 인기 게시물 생성 완료")

    _build_posting_timing(prs, data)
    logger.debug("슬라이드 14-15: 게시 빈도 생성 완료")

    _build_hashtags(prs, data)
    logger.debug("슬라이드 16-17: 해시태그 생성 완료")

    _build_closing(prs, data)
    logger.debug("슬라이드 18: 마무리 생성 완료")

    _build_appendix(prs)
    logger.debug("슬라이드 19: 부록 생성 완료")

    # 저장
    report_dir = data_dir / "report"
    report_dir.mkdir(parents=True, exist_ok=True)
    output_path = report_dir / "report.pptx"
    prs.save(str(output_path))
    logger.info("보고서 저장 완료 → %s", output_path)

    return output_path
